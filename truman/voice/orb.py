"""
orb.py — Truman's visual presence + browser audio bridge
Flask on :5001 serves:
  /         → orb UI (canvas animation)
  /state    → current visual state (polled by UI)
  /audio    → WebSocket carrying 24kHz int16 PCM both ways

The browser handles mic capture + speaker playback. getUserMedia(echoCancellation:true)
gives us production-grade AEC so Truman doesn't hear himself.

States: idle | listening | thinking | speaking
"""

import json
import os
import subprocess
import sys
import threading
import webbrowser
from flask import Flask, jsonify, send_from_directory, request
from flask_sock import Sock

STATIC_DIR = os.path.join(os.path.dirname(__file__), "static")

ORB_URL = "http://localhost:5001"

# ── State constants ───────────────────────────────────────────────────────────
IDLE      = "idle"
LISTENING = "listening"
THINKING  = "thinking"
SPEAKING  = "speaking"

# ── Shared state ──────────────────────────────────────────────────────────────
_state      = IDLE
_state_lock = threading.Lock()

# ── Single-client audio guard ─────────────────────────────────────────────────
# Only ONE browser tab can own the audio stream at a time. If a second tab
# connects, it evicts the first. Without this, multiple tabs split the audio
# queue and you hear fragmented "2-3 bot" echo.
_active_ws          = None
_active_ws_stop_ev  = None
_active_ws_lock     = threading.Lock()

app  = Flask(__name__)
sock = Sock(app)


def set_state(state: str):
    global _state
    with _state_lock:
        _state = state


def get_state() -> str:
    with _state_lock:
        return _state


# ── Flask routes ──────────────────────────────────────────────────────────────
@app.route("/health")
def health():
    return jsonify({"status": "ok", "service": "truman", "mac": "connected" if _mac_ws else "disconnected"})


@app.route("/state")
def state_endpoint():
    return jsonify({"state": get_state()})


@app.route("/")
def index():
    return send_from_directory(STATIC_DIR, "orb.html")


@app.route("/dashboard")
def dashboard():
    return send_from_directory(STATIC_DIR, "dashboard.html")


# ── PWA static assets (Phase 14) ─────────────────────────────────────────────
@app.route("/static/<path:filename>")
def static_files(filename):
    """Serve static assets: icons, manifest, etc."""
    return send_from_directory(STATIC_DIR, filename)


@app.route("/sw.js")
def service_worker():
    """Service worker must be served from root scope to control /dashboard."""
    resp = send_from_directory(STATIC_DIR, "sw.js")
    resp.headers["Service-Worker-Allowed"] = "/"
    resp.headers["Content-Type"] = "application/javascript"
    resp.headers["Cache-Control"] = "no-cache"
    return resp


@app.route("/manifest.json")
def manifest():
    return send_from_directory(STATIC_DIR, "manifest.json")


@app.route("/logs")
def logs_page():
    try:
        from truman.storage import db
        turns = db.recent_turns(50)
        rows = "".join(
            f'<div class="row {"user" if t["role"]=="user" else "bot"}">'
            f'<span class="who">{"Om" if t["role"]=="user" else "Truman"}</span>'
            f'<span class="ts">{(t.get("ts") or "")[:16]}</span>'
            f'<span class="txt">{t["content"][:300]}</span></div>'
            for t in turns
        )
        return f"""<!DOCTYPE html><html><head><meta charset="UTF-8">
<meta name="viewport" content="width=device-width,initial-scale=1">
<title>Truman Logs</title>
<style>
body{{background:#0a0a0f;color:#e8e8f0;font-family:monospace;padding:16px;font-size:13px;}}
h2{{color:#a78bfa;margin-bottom:16px;}}
.row{{padding:8px 0;border-bottom:1px solid #1a1a2e;display:flex;gap:10px;flex-wrap:wrap;}}
.who{{color:#4f46e5;min-width:50px;font-weight:bold;}}
.who.bot{{color:#a78bfa;}}
.ts{{color:#444460;font-size:11px;min-width:100px;}}
.txt{{color:#c8c8e0;flex:1;}}
.user .who{{color:#4f46e5;}}
a{{color:#a78bfa;text-decoration:none;}}
</style></head><body>
<h2>TRUMAN LOGS</h2>
<a href="/dashboard">← dashboard</a>&nbsp;&nbsp;<a href="/health">health</a>
<br><br>{rows or "<p style='color:#444'>no turns yet</p>"}
</body></html>"""
    except Exception as e:
        return f"<pre>error: {e}</pre>", 500


_EMOTIONAL_KW = {"i feel","i'm","im ","i am","i can't","i cant","she ","he ","my ","i've","ive ","i was","i don't","i dont","i need","i want","i hate","i love","honestly","tbh","ngl","bro","man ","struggling","hard","frustrated","anxious","stressed","worried","distracted","sad","lonely","angry","pissed","hurt","confused","lost","scared"}

def _auto_extract_facts(user_input: str) -> None:
    """Background: if message is personal/emotional, extract facts via fast LLM and save."""
    try:
        if len(user_input) < 60:
            return
        lower = user_input.lower()
        if not any(kw in lower for kw in _EMOTIONAL_KW):
            return
        from langchain_core.messages import HumanMessage as _HM
        from truman.core.model_router import run_with_pool
        from truman.storage.db import save_fact, get_all_facts
        existing = [f["fact"].lower() for f in get_all_facts()]
        prompt = (
            "Extract 1-3 short factual statements about the person from this message. "
            "Only extract clear personal facts (feelings, situations, relationships, struggles). "
            "Return ONLY a plain list, one fact per line, no bullets, no numbers, no extra text.\n\n"
            f"Message: {user_input[:600]}"
        )
        result = run_with_pool([_HM(content=prompt)], pool="fast")
        if not result or not result.get("content"):
            return
        for line in result["content"].strip().split("\n"):
            fact = line.strip().lstrip("-•·123456789. ").strip()
            if len(fact) < 10 or len(fact) > 250:
                continue
            # skip if too similar to existing
            if any(fact.lower()[:40] in ex for ex in existing):
                continue
            save_fact(fact, importance=3, source="auto")
    except Exception:
        pass


def _get_session_attachment_summary(session_id: str) -> list:
    """Return active sticky attachment metadata for the dashboard context tray."""
    try:
        from truman.multimodal.session_state import get_session_summary
        return get_session_summary(session_id)
    except Exception:
        return []


def _parse_multimodal_input(raw: str):
    """
    Extract image attach_ids and return clean user_input.
    Images: attach_id extracted, marker stripped (bytes sent live to LLM via content blocks).
    Files (PDF/DOCX etc): attach_id extracted, text content kept inline for LLM.
    Returns (image_ids: list[str], clean_input: str)
    """
    import re as _re
    _attach_pat = _re.compile(r'\[(Image|File):\s*([^|\]]+?)\|attach:([a-f0-9]+)\]', _re.I)
    parts = raw.split('\n---\n') if '\n---\n' in raw else [raw]
    image_ids = []
    clean_parts = []
    for part in parts:
        part = part.strip()
        m = _attach_pat.search(part)
        if m:
            kind = m.group(1).lower()
            aid  = m.group(3)
            if kind == 'image':
                image_ids.append(aid)
                # drop image description — live bytes sent via content blocks
            else:
                # file: keep extracted text, strip just the marker tag
                text_after = part[m.end():].strip()
                fname = m.group(2).strip()
                if text_after:
                    clean_parts.append(f"[File: {fname}]\n{text_after}")
        else:
            if part:
                clean_parts.append(part)
    clean_input = '\n\n'.join(clean_parts).strip() or 'analyze this'
    return image_ids, clean_input


@app.route("/api/chat", methods=["POST"])
def api_chat():
    from flask import request
    import concurrent.futures
    import re as _re
    data = request.get_json(silent=True) or {}
    user_input = (data.get("message") or "").strip()
    session_id = (data.get("session_id") or "default").strip()
    pool_hint  = (data.get("pool") or None)
    if not user_input:
        return jsonify({"error": "empty message"}), 400

    # ── Multimodal: extract image attach_ids, clean user input ───────────────
    image_ids, user_input = _parse_multimodal_input(user_input)

    # Sticky attachment layer (L4)
    try:
        from truman.multimodal.session_state import (
            register_attachments, get_sticky_ids, process_commands
        )
        # Handle "look again" / "drop file" commands
        cmd_ack = process_commands(session_id, user_input)
        # Register any fresh uploads into sticky store
        if image_ids:
            register_attachments(session_id, image_ids)
        # Merge fresh ids with still-active sticky ids (dedup, fresh ids first)
        sticky = get_sticky_ids(session_id)
        all_ids = image_ids + [x for x in sticky if x not in image_ids]
    except Exception as _se:
        print(f"[orb] session_state error: {_se}")
        all_ids = image_ids

    if all_ids and not pool_hint:
        pool_hint = "vision"

    try:
        from truman.text.agent import run as agent_run
        from truman.storage import db as _db

        with concurrent.futures.ThreadPoolExecutor(max_workers=1) as ex:
            future = ex.submit(agent_run, user_input, "", pool_hint, session_id, all_ids)
            try:
                result = future.result(timeout=45)
            except concurrent.futures.TimeoutError:
                return jsonify({"error": "response timed out — try again"}), 504

        # Tick sticky attachment TTL (L4)
        try:
            from truman.multimodal.session_state import tick_turn
            tick_turn(session_id)
        except Exception:
            pass

        # log to SQLite — map browser UUID → SQLite integer session id
        try:
            sid = _db.get_or_create_session(session_id)
            _db.set_session_first_message(session_id, user_input)
            _db.log_turn(sid, "user", user_input)
            _db.log_turn(sid, "assistant", result["response"])
        except Exception:
            pass

        # broadcast to all connected clients (multi-device sync)
        try:
            from truman.storage.notifications import push_turn
            push_turn("user", user_input, session_id)
            push_turn("assistant", result["response"], session_id, {
                "model": result["model"], "pool": result["pool"],
                "tool_calls": result["tool_calls"]
            })
        except Exception:
            pass

        # auto-extract personal facts in background (zero latency impact)
        threading.Thread(
            target=_auto_extract_facts,
            args=(user_input,),
            daemon=True
        ).start()

        mac_status = "connected" if _mac_ws else "disconnected"
        # Prepend drop-file ack if a command was processed
        response_text = result["response"]
        try:
            if cmd_ack:
                response_text = cmd_ack + "\n" + response_text
        except NameError:
            pass
        return jsonify({
            "response":     response_text,
            "model":        result["model"],
            "pool":         result["pool"],
            "tool_calls":   result["tool_calls"],
            "mood":         result["mood"],
            "warnings":     result.get("warnings", []),
            "skill":        result.get("skill", ""),
            "mac":          mac_status,
            "attachments":  _get_session_attachment_summary(session_id),
        })
    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({"error": str(e)}), 500


@app.route("/api/logs")
def api_logs():
    from truman.text.agent import get_error_log
    return jsonify({"logs": get_error_log()})


@app.route("/api/tasks")
def api_tasks():
    """Live status of background tasks (currently: repo ingests).
    Dashboard polls this every 2s while tasks are active."""
    try:
        from truman.storage import db
        tasks = db.active_repo_tasks()
        return jsonify({"tasks": [
            {
                "kind":     "repo",
                "name":     t["name"],
                "url":      t["url"],
                "status":   t["status"],
                "stage":    t.get("stage") or "",
                "progress": t.get("progress") or 0,
                "total":    t.get("total") or 0,
                "files":    t.get("file_count") or 0,
                "error":    t.get("error") or "",
            }
            for t in tasks
        ]})
    except Exception as e:
        return jsonify({"tasks": [], "error": str(e)})


@app.route("/api/power", methods=["GET", "POST"])
def api_power():
    """Om's master kill switch. GET=status, POST=toggle. Truman has no tool for this."""
    from truman.storage.db import killswitch_active, killswitch_set
    if request.method == "GET":
        return jsonify({"on": not killswitch_active()})
    # POST — toggle
    currently_off = killswitch_active()
    killswitch_set(off=currently_off)  # flip
    return jsonify({"on": currently_off})  # was off → now on, and vice versa


@app.route("/api/stream")
def api_stream():
    """SSE endpoint — browser subscribes once, server pushes notifications (task done, errors, etc.)."""
    from flask import Response, stream_with_context
    from truman.storage import notifications as _notif
    import json as _json

    def generate():
        q = _notif.subscribe()
        try:
            yield "data: {\"type\":\"connected\"}\n\n"
            while True:
                try:
                    payload = q.get(timeout=25)
                    yield f"data: {_json.dumps(payload)}\n\n"
                except Exception:
                    yield ": keep-alive\n\n"
        finally:
            _notif.unsubscribe(q)

    return Response(
        stream_with_context(generate()),
        mimetype="text/event-stream",
        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"},
    )


@app.route("/api/trace")
def api_trace():
    """Return persisted trace history (brain node events) for the activity panel."""
    from flask import request as freq
    try:
        from truman.storage.db import get_trace_history
        session_id = freq.args.get("session_id")
        limit = int(freq.args.get("limit", 300))
        rows = get_trace_history(session_id=session_id, limit=limit)
        return jsonify(rows)
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/events")
def api_events():
    """Return last 100 events from DB events table (persisted ring buffer)."""
    from flask import request as freq
    try:
        from truman.storage import db
        kind = freq.args.get("kind")
        date = freq.args.get("date")
        events = db.get_events(limit=100, kind=kind, date=date)
        return jsonify({"events": events})
    except Exception as e:
        return jsonify({"events": [], "error": str(e)})


@app.route("/api/sessions", methods=["GET"])
def api_sessions():
    """All sessions grouped by day for the sidebar."""
    try:
        from truman.storage import db
        rows = db.get_sessions_by_day()
        # group by date
        from collections import OrderedDict
        from datetime import datetime as _dt
        groups = OrderedDict()
        today = _dt.now().date()
        for r in rows:
            try:
                d = _dt.fromisoformat(r["started_at"]).date()
            except Exception:
                d = today
            delta = (today - d).days
            if delta == 0:
                label = "Today"
            elif delta == 1:
                label = "Yesterday"
            else:
                label = d.strftime("%B %-d")
            groups.setdefault(label, []).append(r)
        return jsonify({"groups": [{"day": k, "sessions": v} for k, v in groups.items()]})
    except Exception as e:
        return jsonify({"groups": [], "error": str(e)})


@app.route("/api/sessions/<browser_id>", methods=["PATCH"])
def api_session_rename(browser_id):
    try:
        from truman.storage import db
        from flask import request as freq
        label = (freq.get_json(force=True) or {}).get("label", "").strip()
        if label:
            db.update_session_label(browser_id, label[:50])
        return jsonify({"ok": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/sessions/<browser_id>", methods=["DELETE"])
def api_session_delete(browser_id):
    try:
        from truman.storage import db
        db.delete_session(browser_id)
        return jsonify({"ok": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/history")
def api_history():
    """Return turns from SQLite. If session_id passed, return only that session's turns."""
    try:
        from flask import request as freq
        from truman.storage import db
        sid = freq.args.get("session_id")
        turns = db.session_turns(sid) if sid else db.recent_turns(30)
        return jsonify({"turns": turns})
    except Exception as e:
        return jsonify({"turns": [], "error": str(e)})


# ── User Facts ───────────────────────────────────────────────────────────────
@app.route("/api/facts", methods=["GET"])
def api_facts_get():
    try:
        from truman.storage import db
        return jsonify({"facts": db.get_all_facts()})
    except Exception as e:
        return jsonify({"facts": [], "error": str(e)})

@app.route("/api/facts", methods=["POST"])
def api_facts_post():
    try:
        from truman.storage import db
        from flask import request as freq
        data = freq.get_json(force=True)
        fact = (data.get("fact") or "").strip()
        if not fact:
            return jsonify({"error": "empty fact"}), 400
        importance = int(data.get("importance", 3))
        source = data.get("source", "manual")
        fid = db.save_fact(fact, importance, source)
        return jsonify({"id": fid, "ok": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/api/facts/<int:fact_id>", methods=["DELETE"])
def api_facts_delete(fact_id):
    try:
        from truman.storage import db
        db.delete_fact(fact_id)
        return jsonify({"ok": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


# ── Persona Rules (Phase 13) ─────────────────────────────────────────────────
@app.route("/api/rules", methods=["GET"])
def api_rules_get():
    try:
        from truman.storage import db
        return jsonify({"rules": db.get_all_rules()})
    except Exception as e:
        return jsonify({"rules": [], "error": str(e)})

@app.route("/api/rules", methods=["POST"])
def api_rules_post():
    try:
        from truman.storage import db
        from flask import request as freq
        data = freq.get_json(force=True)
        rule = (data.get("rule") or "").strip()
        if not rule:
            return jsonify({"error": "empty rule"}), 400
        rid = db.add_rule(rule, source=data.get("source", "manual"))
        return jsonify({"id": rid, "ok": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/api/rules/<int:rule_id>", methods=["PATCH"])
def api_rules_toggle(rule_id):
    try:
        from truman.storage import db
        from flask import request as freq
        data = freq.get_json(force=True)
        db.toggle_rule(rule_id, int(data.get("active", 1)))
        return jsonify({"ok": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/api/rules/<int:rule_id>", methods=["DELETE"])
def api_rules_delete(rule_id):
    try:
        from truman.storage import db
        db.delete_rule(rule_id)
        return jsonify({"ok": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


# ── Reply contacts (Phase 15D) ───────────────────────────────────────────────
@app.route("/api/contacts", methods=["GET"])
def api_contacts_get():
    try:
        from truman.storage import db
        return jsonify({"contacts": db.list_reply_contacts()})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/api/contacts", methods=["POST"])
def api_contacts_add():
    try:
        from truman.storage import db
        name = (request.json or {}).get("name", "").strip()
        if not name:
            return jsonify({"error": "name required"}), 400
        added = db.add_reply_contact(name)
        return jsonify({"ok": True, "added": added})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/api/contacts/<int:contact_id>", methods=["DELETE"])
def api_contacts_delete(contact_id):
    try:
        from truman.storage import db
        with db._conn() as c:
            c.execute("DELETE FROM reply_contacts WHERE id = ?", (contact_id,))
        return jsonify({"ok": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


# ── Web Push (Phase 14) ───────────────────────────────────────────────────────
@app.route("/api/push/vapid-public-key", methods=["GET"])
def api_push_vapid_key():
    """Return VAPID public key so frontend can subscribe."""
    try:
        from truman.delivery.web_push import get_public_key
        return jsonify({"key": get_public_key()})
    except Exception as e:
        return jsonify({"key": "", "error": str(e)})

@app.route("/api/push/subscribe", methods=["POST"])
def api_push_subscribe():
    """Save a push subscription from the browser."""
    try:
        from truman.storage import db
        data = request.get_json(force=True)
        endpoint = data.get("endpoint", "").strip()
        p256dh   = (data.get("keys") or {}).get("p256dh", "").strip()
        auth     = (data.get("keys") or {}).get("auth", "").strip()
        if not endpoint or not p256dh or not auth:
            return jsonify({"error": "missing subscription fields"}), 400
        db.save_push_sub(endpoint, p256dh, auth)
        return jsonify({"ok": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/api/push/unsubscribe", methods=["POST"])
def api_push_unsubscribe():
    try:
        from truman.storage import db
        data = request.get_json(force=True)
        endpoint = (data.get("endpoint") or "").strip()
        if endpoint:
            db.delete_push_sub(endpoint)
        return jsonify({"ok": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


# ── File Upload ───────────────────────────────────────────────────────────────
TEXT_EXTS = {".py",".js",".ts",".html",".css",".json",".md",".txt",".csv",".xml",".yaml",".yml",".sh",".env"}
IMAGE_EXTS = {".png",".jpg",".jpeg",".gif",".webp",".bmp",".tiff"}

def _extract_text(filename: str, data: bytes) -> str:
    import io
    ext = os.path.splitext(filename)[1].lower()
    if ext in TEXT_EXTS:
        return data.decode("utf-8", errors="replace")
    if ext == ".pdf":
        import pdfplumber
        with pdfplumber.open(io.BytesIO(data)) as pdf:
            return "\n".join(p.extract_text() or "" for p in pdf.pages)
    if ext in (".docx",):
        from docx import Document
        return "\n".join(p.text for p in Document(io.BytesIO(data)).paragraphs)
    if ext in (".xlsx",):
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        rows = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                rows.append("\t".join("" if v is None else str(v) for v in row))
        return "\n".join(rows)
    if ext in (".pptx",):
        from pptx import Presentation
        lines = []
        for slide in Presentation(io.BytesIO(data)).slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    lines.append(shape.text_frame.text)
        return "\n".join(lines)
    if ext in IMAGE_EXTS:
        return "__IMAGE__"   # signal to handler to use vision model
    return data.decode("utf-8", errors="replace")


@app.route("/api/upload", methods=["POST"])
def api_upload():
    from flask import request
    import uuid as _uuid
    f = request.files.get("file")
    if not f:
        return jsonify({"error": "no file"}), 400
    data = f.read()
    ext  = os.path.splitext(f.filename)[1].lower()
    mime_map = {".png":"image/png",".jpg":"image/jpeg",".jpeg":"image/jpeg",
                ".gif":"image/gif",".webp":"image/webp",".bmp":"image/bmp",
                ".pdf":"application/pdf",".txt":"text/plain",".md":"text/plain",
                ".docx":"application/vnd.openxmlformats-officedocument.wordprocessingml.document"}
    mime = mime_map.get(ext, "application/octet-stream")

    # ── Persist raw bytes to SQLite (Phase 15 — survives refresh forever) ──
    attach_id = _uuid.uuid4().hex[:16]
    try:
        from truman.storage.db import save_attachment
        save_attachment(attach_id, f.filename, mime, data)
    except Exception as e:
        print(f"[Upload] attachment save failed: {e}")
        attach_id = None

    try:
        text = _extract_text(f.filename, data)
    except Exception as e:
        return jsonify({"error": f"parse failed: {e}"}), 500

    # Images: bytes are stored above — no describe-once. LLM sees live bytes via content blocks.
    if text == "__IMAGE__":
        text = ""

    if len(text) > 30000:
        text = text[:30000] + "\n\n[...truncated — document exceeds 30K chars]"

    return jsonify({"filename": f.filename, "text": text,
                    "attach_id": attach_id, "mime": mime})


@app.route("/api/attachments/<attach_id>")
def api_attachment(attach_id):
    """Serve a stored file/image by attach_id. Used by dashboard to render persistent images."""
    try:
        from flask import request as freq, Response
        from truman.storage.db import get_attachment
        att = get_attachment(attach_id)
        if not att:
            return jsonify({"error": "not found"}), 404
        # inline for images, attachment for everything else
        disposition = "inline" if att["mime_type"].startswith("image/") else "attachment"
        return Response(
            att["data"],
            mimetype=att["mime_type"],
            headers={"Content-Disposition": f'{disposition}; filename="{att["filename"]}"',
                     "Cache-Control": "public, max-age=31536000"}  # cache 1yr — content never changes
        )
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/attachments/session/<session_id>")
def api_session_attachments(session_id):
    """Return active sticky attachments for a session (for dashboard context tray)."""
    return jsonify({"attachments": _get_session_attachment_summary(session_id)})


@app.route("/api/attachments/session/<session_id>/drop", methods=["POST"])
def api_drop_attachments(session_id):
    """Drop sticky attachments for a session. Body: {"kind": "all"|"image"|"file"}"""
    from flask import request as _req
    kind = (_req.get_json(silent=True) or {}).get("kind", "all")
    try:
        from truman.multimodal.session_state import clear_attachments
        n = clear_attachments(session_id, kind)
        return jsonify({"dropped": n})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


# ── Phase 15: Boss / WhatsApp intake ─────────────────────────────────────────

@app.route("/api/boss_message", methods=["POST"])
def api_boss_message():
    """
    iPhone Shortcut forwards any WhatsApp/iMessage → here.
    Body: {"from": "Adam", "text": "...", "source": "whatsapp"}
    Truman drafts reply → Telegram approval flow.
    """
    from flask import request as freq
    data   = freq.get_json(silent=True) or {}
    sender = (data.get("from") or data.get("sender") or "Unknown").strip()
    text   = (data.get("text") or data.get("message") or "").strip()
    source = (data.get("source") or "whatsapp").strip()
    if not text:
        return jsonify({"error": "empty message"}), 400
    try:
        from truman.integrations.boss_handler import handle_incoming
        result = handle_incoming(sender, text, source)
        return jsonify(result)
    except Exception as e:
        return jsonify({"error": str(e)}), 500


# ── Mac Bridge WebSocket (Railway side) ──────────────────────────────────────
BRIDGE_SECRET = os.environ.get("BRIDGE_SECRET", "truman-bridge-secret")
_mac_ws = None
_mac_ws_lock = threading.Lock()
_pending: dict = {}   # id → asyncio.Event (not needed — flask-sock is sync)
_pending_results: dict = {}


@sock.route("/mac-bridge")
def mac_bridge_ws(ws):
    """Accepts persistent connection from the Mac Bridge daemon."""
    global _mac_ws
    secret = ws.environ.get("HTTP_X_BRIDGE_SECRET", "")
    if secret != BRIDGE_SECRET:
        ws.close()
        return
    with _mac_ws_lock:
        _mac_ws = ws
    print("[Bridge] Mac connected.")
    try:
        while True:
            raw = ws.receive()
            if raw is None:
                break
            try:
                msg = json.loads(raw)
                req_id = msg.get("id", "")
                if req_id in _pending:
                    _pending_results[req_id] = msg
                    _pending.pop(req_id)
            except Exception:
                pass
    finally:
        with _mac_ws_lock:
            if _mac_ws is ws:
                _mac_ws = None
        print("[Bridge] Mac disconnected.")


def mac_request(action: str, payload: dict, timeout: float = 15.0) -> dict:
    """Send a request to the Mac Bridge and wait for response."""
    import uuid, time as _time
    with _mac_ws_lock:
        ws = _mac_ws
    if not ws:
        return {"ok": False, "error": "Mac bridge not connected — laptop may be asleep."}
    req_id = str(uuid.uuid4())[:8]
    _pending[req_id] = True
    try:
        ws.send(json.dumps({"id": req_id, "action": action, **payload}))
    except Exception as e:
        _pending.pop(req_id, None)
        return {"ok": False, "error": str(e)}
    deadline = _time.time() + timeout
    while _time.time() < deadline:
        if req_id in _pending_results:
            return _pending_results.pop(req_id)
        _time.sleep(0.1)
    _pending.pop(req_id, None)
    return {"ok": False, "error": "Mac bridge timed out."}


# ── Audio WebSocket ───────────────────────────────────────────────────────────
@sock.route("/audio")
def audio_ws(ws):
    """Bridge browser ⇄ realtime.py queues.

    Inbound  (browser binary frames, 24kHz int16 PCM) → realtime.mic_in
    Outbound (realtime.audio_out frames)              → browser binary frames
    A None sentinel on audio_out is translated to a JSON flush message.

    Single-client: when a new tab connects, the previous one is evicted so
    only ONE browser pulls from the audio queue. Prevents "multiple bots
    speaking" echo when old tabs linger across restarts.
    """
    from truman.voice import realtime  # deferred import to avoid circular dep

    # Auto-start Realtime session on Railway (no hotkey available in cloud).
    # If the session was just closing (_session_active briefly True while finalizer runs),
    # wait up to 3s for it to fully close before starting a new one.
    if realtime._session_active:
        # session appears active — could be live OR mid-shutdown.
        # Give it up to 3s to settle; if still active after, assume live (reuse it).
        import time as _t
        deadline = _t.time() + 3.0
        while realtime._session_active and _t.time() < deadline:
            _t.sleep(0.1)
    if not realtime._session_active:
        realtime.start_session()

    global _active_ws, _active_ws_stop_ev
    stop_flag = threading.Event()

    # Evict any prior client before we take over the audio queue
    with _active_ws_lock:
        if _active_ws_stop_ev is not None:
            _active_ws_stop_ev.set()
            try:
                # Tell the old tab to shut itself (prevents auto-reconnect storm)
                _active_ws.send(json.dumps({"type": "evicted"}))
                _active_ws.close()
            except Exception:
                pass
        _active_ws         = ws
        _active_ws_stop_ev = stop_flag
        print("[Orb] Audio client connected (sole consumer)")

    def reader():
        # Browser → OpenAI
        try:
            while not stop_flag.is_set():
                data = ws.receive(timeout=1)
                if data is None:
                    continue
                if isinstance(data, (bytes, bytearray)):
                    try:
                        realtime.mic_in.put_nowait(bytes(data))
                    except Exception:
                        pass   # drop on backpressure
                # text frames (control messages) ignored for now
        except Exception:
            pass
        finally:
            stop_flag.set()

    t = threading.Thread(target=reader, daemon=True)
    t.start()

    # OpenAI → Browser
    try:
        while not stop_flag.is_set():
            try:
                frame = realtime.audio_out.get(timeout=0.1)
            except Exception:
                if not t.is_alive():
                    break
                continue

            try:
                if frame is None:
                    ws.send(json.dumps({"type": "flush"}))
                elif isinstance(frame, dict):
                    ws.send(json.dumps(frame))   # transcript / control messages
                else:
                    ws.send(frame)               # binary audio
            except Exception:
                break
    finally:
        stop_flag.set()
        with _active_ws_lock:
            if _active_ws is ws:
                _active_ws         = None
                _active_ws_stop_ev = None


# ── Public API ────────────────────────────────────────────────────────────────
def _open_browser():
    """Open the orb UI in a NEW browser window and bring it to the front.
    Uses osascript on Mac so the window can't hide behind existing tabs/spaces."""
    try:
        if sys.platform == "darwin":
            # osascript: open URL in a brand-new Chrome/Safari/default window AND activate it
            script = f'''
            tell application "System Events"
                set frontApp to name of (first application process whose frontmost is true)
            end tell
            try
                tell application "Google Chrome"
                    activate
                    make new window
                    set URL of active tab of front window to "{ORB_URL}"
                end tell
            on error
                try
                    tell application "Safari"
                        activate
                        make new document with properties {{URL:"{ORB_URL}"}}
                    end tell
                on error
                    do shell script "open '{ORB_URL}'"
                end try
            end try
            '''
            subprocess.Popen(["osascript", "-e", script])
        else:
            webbrowser.open_new(ORB_URL)
        print(f"[Orb] Opened {ORB_URL} in new window")
    except Exception as e:
        print(f"[Orb] Could not auto-open browser ({e}). Visit {ORB_URL} manually.")


def run(port: int = 5001, open_browser: bool = True):
    """Start orb server in background thread. Non-blocking."""
    t = threading.Thread(target=_serve, args=(port,), daemon=True)
    t.start()
    if open_browser:
        threading.Timer(1.5, _open_browser).start()


def start():
    run()


def stop():
    pass   # daemon thread dies with the process


def _serve(port: int = 5001):
    import logging
    log = logging.getLogger("werkzeug")
    log.setLevel(logging.ERROR)
    host = "0.0.0.0" if os.environ.get("RAILWAY_ENVIRONMENT") else "127.0.0.1"
    app.run(host=host, port=port, debug=False, use_reloader=False, threaded=True)


# HTML served from static/ — see truman/voice/static/orb.html and dashboard.html
